from __future__ import annotations

import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


PROJECT_DIR = Path(__file__).resolve().parents[1]
WORKSPACE_DIR = PROJECT_DIR.parent
DOCS_DIR = WORKSPACE_DIR / "docs"
if str(PROJECT_DIR) not in sys.path:
    sys.path.insert(0, str(PROJECT_DIR))

from vkr_classifier.config import get_settings  # noqa: E402


OUTPUT_NAME = "Презентация ВКР.pptx"
TEMPLATE_NAME = "Шаблон презентации.pptx"
TITLE_COLOR = RGBColor(47, 48, 131)
TEXT_COLOR = RGBColor(36, 39, 49)
ACCENT = RGBColor(83, 104, 215)
CARD_FILL = RGBColor(247, 249, 255)
CARD_ALT = RGBColor(232, 240, 255)


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slide = slide_id_list[index]
    rel_id = slide.rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide)


def set_text_run_style(run, *, size: int, color: RGBColor = TEXT_COLOR, bold: bool = False) -> None:
    run.font.name = "Montserrat"
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold


def set_title(shape, text: str, font_size: int = 28) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    for run in paragraph.runs:
        set_text_run_style(run, size=font_size, color=TITLE_COLOR, bold=True)


def set_body(
    shape,
    lines: list[str],
    *,
    font_size: int = 20,
    color: RGBColor = TEXT_COLOR,
    bold_first: bool = False,
) -> None:
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.level = 0
        paragraph.space_after = Pt(6)
        for run in paragraph.runs:
            set_text_run_style(run, size=font_size, color=color, bold=bold_first and index == 0)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list[str],
    *,
    font_size: int = 20,
    color: RGBColor = TEXT_COLOR,
    bold_first: bool = False,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    set_body(textbox, lines, font_size=font_size, color=color, bold_first=bold_first)
    return textbox


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title: str,
    body_lines: list[str],
    *,
    fill_color: RGBColor = CARD_FILL,
    title_size: int = 17,
    body_size: int = 14,
) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1.25)

    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(10)
    text_frame.margin_right = Pt(10)
    text_frame.margin_top = Pt(10)
    text_frame.margin_bottom = Pt(8)

    title_paragraph = text_frame.paragraphs[0]
    title_paragraph.text = title
    title_paragraph.space_after = Pt(6)
    for run in title_paragraph.runs:
        set_text_run_style(run, size=title_size, color=TITLE_COLOR, bold=True)

    for line in body_lines:
        paragraph = text_frame.add_paragraph()
        paragraph.text = line
        paragraph.space_after = Pt(3)
        for run in paragraph.runs:
            set_text_run_style(run, size=body_size)


def add_stat_card(
    slide,
    left,
    top,
    width,
    height,
    title: str,
    value_lines: list[str],
    *,
    fill_color: RGBColor = CARD_ALT,
) -> None:
    add_card(
        slide,
        left,
        top,
        width,
        height,
        title,
        value_lines,
        fill_color=fill_color,
        title_size=18,
        body_size=16,
    )


def clear_content_placeholders(slide) -> None:
    for index in (0, 2):
        if index < len(slide.shapes):
            shape = slide.shapes[index]
            if hasattr(shape, "text_frame"):
                shape.text_frame.clear()


def set_slide_number(shape, number: int) -> None:
    set_body(shape, [str(number)], font_size=14, color=TITLE_COLOR)


def build_presentation() -> Path:
    settings = get_settings(PROJECT_DIR)
    summary = pd.read_csv(settings.summary_metrics_path)
    text_metrics = summary[summary["model"] == "Текстовая модель"].iloc[0]
    image_metrics = summary[summary["model"] == "Модель изображений"].iloc[0]

    presentation = Presentation(DOCS_DIR / TEMPLATE_NAME)

    for index in range(len(presentation.slides) - 1, 8, -1):
        remove_slide(presentation, index)
    remove_slide(presentation, 0)

    title_slide = presentation.slides[0]
    set_title(
        title_slide.shapes[3],
        "Разработка системы интеллектуальной классификации и пакетной сортировки документов на языке Python с low-code интерфейсом",
        font_size=23,
    )
    set_body(title_slide.shapes[1], ["Студент: Медведев Илья", "Группа: не указана"], font_size=18)
    set_body(title_slide.shapes[0], ["Научный руководитель: данные уточняются"], font_size=16)

    relevance_slide = presentation.slides[1]
    clear_content_placeholders(relevance_slide)
    set_title(relevance_slide.shapes[1], "Практическая задача и актуальность")
    add_card(
        relevance_slide,
        Inches(0.72),
        Inches(1.55),
        Inches(2.85),
        Inches(1.5),
        "Смешанные архивы",
        ["Организации получают документы текстом, сканами и пакетами файлов."],
        fill_color=CARD_FILL,
        title_size=16,
        body_size=13,
    )
    add_card(
        relevance_slide,
        Inches(3.72),
        Inches(1.55),
        Inches(2.9),
        Inches(1.5),
        "Ручная сортировка",
        ["Разбор массива документов занимает время и плохо масштабируется."],
        fill_color=CARD_ALT,
        title_size=16,
        body_size=13,
    )
    add_card(
        relevance_slide,
        Inches(6.77),
        Inches(1.55),
        Inches(2.9),
        Inches(1.5),
        "ML-подход",
        ["Тип документа можно определять по тексту и по структуре макета."],
        fill_color=CARD_FILL,
        title_size=16,
        body_size=13,
    )
    add_card(
        relevance_slide,
        Inches(9.82),
        Inches(1.55),
        Inches(2.35),
        Inches(1.5),
        "Low-code UI",
        ["Gradio дает быстрый прикладной интерфейс без отдельного frontend."],
        fill_color=CARD_ALT,
        title_size=15,
        body_size=12,
    )
    add_card(
        relevance_slide,
        Inches(0.95),
        Inches(3.55),
        Inches(11.15),
        Inches(1.45),
        "Цель работы",
        ["Разработать локальное приложение, которое классифицирует документы и автоматически сортирует ZIP-архив по типам: договор, счет, приказ, служебная записка, отчет."],
        fill_color=CARD_ALT,
        title_size=20,
        body_size=16,
    )
    set_body(
        relevance_slide.shapes[2],
        ["Результат ВКР - не абстрактный классификатор, а прикладной сервис первичной маршрутизации документов."],
        font_size=12,
    )
    set_slide_number(relevance_slide.shapes[3], 2)

    tasks_slide = presentation.slides[2]
    clear_content_placeholders(tasks_slide)
    set_title(tasks_slide.shapes[1], "Поддерживаемые сценарии")
    add_card(
        tasks_slide,
        Inches(0.8),
        Inches(1.55),
        Inches(3.7),
        Inches(1.55),
        "1. Текст документа",
        ["Ввод текста и мгновенное определение типа документа."],
        fill_color=CARD_FILL,
    )
    add_card(
        tasks_slide,
        Inches(4.78),
        Inches(1.55),
        Inches(3.7),
        Inches(1.55),
        "2. Скан документа",
        ["Загрузка изображения и классификация по макету страницы."],
        fill_color=CARD_ALT,
    )
    add_card(
        tasks_slide,
        Inches(8.76),
        Inches(1.55),
        Inches(3.0),
        Inches(1.55),
        "3. ZIP-архив",
        ["Пакетная сортировка массива файлов и выдача нового архива."],
        fill_color=CARD_FILL,
    )
    add_card(
        tasks_slide,
        Inches(0.95),
        Inches(3.55),
        Inches(10.9),
        Inches(1.55),
        "Классы документов",
        ["Договор", "Счет", "Приказ", "Служебная записка", "Отчет"],
        fill_color=CARD_ALT,
        title_size=18,
        body_size=16,
    )
    set_body(
        tasks_slide.shapes[2],
        ["Ключевое расширение проекта - пакетная сортировка архива, добавляющая практическую ценность приложению."],
        font_size=12,
    )
    set_slide_number(tasks_slide.shapes[3], 3)

    architecture_slide = presentation.slides[3]
    clear_content_placeholders(architecture_slide)
    set_title(architecture_slide.shapes[1], "Архитектура решения")
    architecture_slide.shapes.add_picture(
        str(settings.architecture_figure),
        left=Inches(0.82),
        top=Inches(1.48),
        width=Inches(11.1),
        height=Inches(3.72),
    )
    add_card(
        architecture_slide,
        Inches(0.85),
        Inches(5.35),
        Inches(3.35),
        Inches(0.82),
        "Интерфейс",
        ["Gradio собирает три прикладных сценария в одном окне."],
        fill_color=CARD_FILL,
        title_size=15,
        body_size=11,
    )
    add_card(
        architecture_slide,
        Inches(4.55),
        Inches(5.35),
        Inches(3.05),
        Inches(0.82),
        "Сервисный контур",
        ["FastAPI и сервисный слой управляют классификацией и логированием."],
        fill_color=CARD_ALT,
        title_size=15,
        body_size=11,
    )
    add_card(
        architecture_slide,
        Inches(7.95),
        Inches(5.35),
        Inches(3.85),
        Inches(0.82),
        "Хранилище и архивы",
        ["SQLite хранит историю, архивный модуль формирует CSV и ZIP-результат."],
        fill_color=CARD_FILL,
        title_size=15,
        body_size=11,
    )
    set_body(
        architecture_slide.shapes[2],
        ["Архитектура ориентирована на переносимость и локальный запуск без внешних сервисов."],
        font_size=12,
    )
    set_slide_number(architecture_slide.shapes[3], 4)

    implementation_slide = presentation.slides[4]
    clear_content_placeholders(implementation_slide)
    set_title(implementation_slide.shapes[1], "Программная реализация")
    implementation_slide.shapes.add_picture(
        str(settings.interaction_figure),
        left=Inches(5.0),
        top=Inches(1.38),
        width=Inches(6.95),
        height=Inches(4.55),
    )
    add_card(
        implementation_slide,
        Inches(0.7),
        Inches(1.65),
        Inches(4.0),
        Inches(1.65),
        "Что реализовано",
        ["текстовая модель TF-IDF + Logistic Regression", "визуальная модель Random Forest", "пакетная сортировка ZIP-архива"],
        fill_color=CARD_FILL,
    )
    add_card(
        implementation_slide,
        Inches(0.7),
        Inches(3.65),
        Inches(4.0),
        Inches(1.62),
        "Сопровождающие блоки",
        ["SQLite-журнал истории", "демо-архив и отчеты", "автотесты и скрипты пересборки"],
        fill_color=CARD_ALT,
    )
    set_body(
        implementation_slide.shapes[2],
        ["Интерфейс показывает не только прогноз, но и метрики моделей и историю пакетных прогонов."],
        font_size=12,
    )
    set_slide_number(implementation_slide.shapes[3], 5)

    results_slide = presentation.slides[5]
    clear_content_placeholders(results_slide)
    set_title(results_slide.shapes[1], "Результаты тестирования и экспериментов")
    results_slide.shapes.add_picture(
        str(settings.model_comparison_figure),
        left=Inches(6.08),
        top=Inches(1.55),
        width=Inches(5.85),
        height=Inches(3.75),
    )
    add_stat_card(
        results_slide,
        Inches(0.72),
        Inches(1.55),
        Inches(4.95),
        Inches(1.45),
        "Текстовая модель",
        [f"Accuracy: {text_metrics['accuracy']:.4f}", f"F1-score: {text_metrics['f1_score']:.4f}"],
        fill_color=CARD_ALT,
    )
    add_stat_card(
        results_slide,
        Inches(0.72),
        Inches(3.1),
        Inches(4.95),
        Inches(1.45),
        "Модель сканов",
        [f"Accuracy: {image_metrics['accuracy']:.4f}", f"F1-score: {image_metrics['f1_score']:.4f}"],
        fill_color=CARD_FILL,
    )
    add_textbox(
        results_slide,
        Inches(0.78),
        Inches(4.9),
        Inches(5.2),
        Inches(0.78),
        ["16 автотестов, покрытие кода выше 85%, интерактивное время ответа после прогрева."],
        font_size=13,
        color=TITLE_COLOR,
        bold_first=True,
    )
    set_body(
        results_slide.shapes[2],
        ["Пакетный режим формирует CSV-сводку и новый ZIP-архив с разложением по типам документов."],
        font_size=12,
    )
    set_slide_number(results_slide.shapes[3], 6)

    value_slide = presentation.slides[6]
    clear_content_placeholders(value_slide)
    set_title(value_slide.shapes[1], "Практическая ценность и развитие")
    value_slide.shapes.add_picture(
        str(settings.use_case_figure),
        left=Inches(7.02),
        top=Inches(1.62),
        width=Inches(4.82),
        height=Inches(3.35),
    )
    add_card(
        value_slide,
        Inches(0.78),
        Inches(1.62),
        Inches(5.8),
        Inches(1.18),
        "Практический эффект",
        ["Сервис выполняет первичную расфасовку архива документов без ручного открытия каждого файла."],
        fill_color=CARD_FILL,
        title_size=18,
        body_size=14,
    )
    add_card(
        value_slide,
        Inches(0.78),
        Inches(2.95),
        Inches(5.8),
        Inches(1.18),
        "Инженерная зрелость",
        ["Есть рабочее приложение, локальная БД, тесты, диаграммы, пояснительная записка и презентация."],
        fill_color=CARD_ALT,
        title_size=18,
        body_size=14,
    )
    add_card(
        value_slide,
        Inches(0.78),
        Inches(4.28),
        Inches(5.8),
        Inches(1.18),
        "Направления развития",
        ["Подключение реальных корпусов документов, OCR и более сложных мультимодальных моделей."],
        fill_color=CARD_FILL,
        title_size=18,
        body_size=14,
    )
    set_body(
        value_slide.shapes[2],
        ["Разработанное ПО может выступать как прототип сервиса первичной маршрутизации корпоративных документов."],
        font_size=12,
    )
    set_slide_number(value_slide.shapes[3], 7)

    thanks_slide = presentation.slides[7]
    set_title(thanks_slide.shapes[3], "Спасибо за внимание", font_size=26)
    set_body(thanks_slide.shapes[1], ["Студент: Медведев Илья", "Группа: не указана"], font_size=18)
    set_body(thanks_slide.shapes[0], ["Научный руководитель: данные уточняются"], font_size=16)

    DOCS_DIR.mkdir(parents=True, exist_ok=True)
    output_path = DOCS_DIR / OUTPUT_NAME
    presentation.save(output_path)
    return output_path


if __name__ == "__main__":
    result = build_presentation()
    print(result)
